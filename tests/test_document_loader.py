"""文档加载器测试：覆盖 TXT/MD/CSV（零依赖）与多格式支持声明。"""
import csv
import tempfile
from pathlib import Path

from src.document_loader import load_document, supported_extensions


def test_supported_extensions_includes_new_formats():
    exts = supported_extensions()
    for expected in ["pdf", "txt", "md", "docx", "pptx", "html", "csv"]:
        assert expected in exts


def test_read_txt_and_md_page_one():
    with tempfile.TemporaryDirectory() as d:
        p = Path(d) / "a.md"
        p.write_text("# 标题\n正文内容", encoding="utf-8")
        pages = load_document(p)
        assert len(pages) == 1
        assert pages[0]["page"] == 1
        assert "标题" in pages[0]["text"]


def test_read_csv_one_page_per_row():
    with tempfile.TemporaryDirectory() as d:
        p = Path(d) / "data.csv"
        with p.open("w", encoding="utf-8", newline="") as f:
            w = csv.writer(f)
            w.writerow(["name", "score"])
            w.writerow(["Alice", "90"])
            w.writerow(["Bob", "85"])
        pages = load_document(p)
        assert len(pages) == 3  # 表头 + 两行
        assert pages[0]["page"] == 1
        assert "name" in pages[0]["text"]
        assert "Alice" in pages[1]["text"]


def test_unsupported_extension_raises():
    with tempfile.TemporaryDirectory() as d:
        p = Path(d) / "x.xyz"
        p.write_text("hi")
        try:
            load_document(p)
            assert False, "应当抛出 ValueError"
        except ValueError:
            pass


def test_docx_and_pptx_and_html_parse_if_available():
    # 这些格式依赖可选库，可用时验证；不可用时跳过。
    with tempfile.TemporaryDirectory() as d:
        # DOCX
        try:
            from docx import Document

            p = Path(d) / "doc.docx"
            doc = Document()
            doc.add_paragraph("第一段")
            doc.add_paragraph("第二段")
            doc.save(str(p))
            pages = load_document(p)
            assert pages and pages[0]["page"] == 1
            assert "第一段" in pages[0]["text"]
        except ImportError:
            pass

        # PPTX
        try:
            from pptx import Presentation

            p = Path(d) / "s.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.add_textbox(0, 0, 100, 100).text_frame.text = "幻灯片文字"
            prs.save(str(p))
            pages = load_document(p)
            assert any("幻灯片文字" in pg["text"] for pg in pages)
        except ImportError:
            pass

        # HTML
        try:
            from bs4 import BeautifulSoup

            p = Path(d) / "page.html"
            p.write_text("<html><body><p>网页正文</p><script>ignore()</script></body></html>", encoding="utf-8")
            pages = load_document(p)
            assert pages and "网页正文" in pages[0]["text"]
            assert "ignore" not in pages[0]["text"]
        except ImportError:
            pass

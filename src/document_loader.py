"""负责解析多种格式文档，并将不同格式文档统一转换为按页组织的文本结构。

格式支持（中期增强：扩展文档格式 + 升级 PDF 解析）：
- PDF   ：优先使用 PyMuPDF(fitz) 提取，质量更高；未安装时回退到 PyPDF。
- DOCX  ：python-docx，整篇作为一个逻辑页（page=1）。
- PPTX  ：python-pptx，每张幻灯片作为一个逻辑页。
- HTML  ：BeautifulSoup 去除标签，整篇作为一个逻辑页。
- CSV   ：标准库 csv，按行解析，每行作为一个逻辑页。
- TXT/MD：纯文本读取，整篇作为一个逻辑页。

统一输出：List[{"text": str, "page": int}]，与后续切块、向量化链路完全兼容。
"""
from pathlib import Path
from typing import Any, Dict, List


# ---------------------------------------------------------------------------
# 解析器
# ---------------------------------------------------------------------------
def read_pdf(file_path: Path) -> List[Dict[str, Any]]:
    """PDF 文本提取：优先 PyMuPDF，回退 PyPDF。"""
    try:
        import fitz  # PyMuPDF

        pages: List[Dict[str, Any]] = []
        with fitz.open(str(file_path)) as doc:
            for page_index, page in enumerate(doc):
                text = page.get_text().strip()
                if text:
                    pages.append({"text": text, "page": page_index + 1})
        return pages
    except ImportError:
        # 回退到 PyPDF
        from pypdf import PdfReader

        pages = []
        reader = PdfReader(str(file_path))
        for page_index, page in enumerate(reader.pages):
            text = (page.extract_text() or "").strip()
            if text:
                pages.append({"text": text, "page": page_index + 1})
        return pages


def read_docx(file_path: Path) -> List[Dict[str, Any]]:
    """DOCX 解析：合并段落为整篇文本，作为一个逻辑页。"""
    import docx

    document = docx.Document(str(file_path))
    paragraphs = [
        para.text.strip()
        for para in document.paragraphs
        if para.text and para.text.strip()
    ]
    text = "\n".join(paragraphs).strip()
    if not text:
        return []
    return [{"text": text, "page": 1}]


def read_pptx(file_path: Path) -> List[Dict[str, Any]]:
    """PPTX 解析：每张幻灯片作为一个逻辑页。"""
    from pptx import Presentation

    presentation = Presentation(str(file_path))
    pages: List[Dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        text = "\n".join(texts).strip()
        if text:
            pages.append({"text": text, "page": index})
    return pages


def read_html(file_path: Path) -> List[Dict[str, Any]]:
    """HTML 解析：去除标签保留文本，作为一个逻辑页。"""
    from bs4 import BeautifulSoup

    raw = file_path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    text = soup.get_text(separator="\n").strip()
    if not text:
        return []
    return [{"text": text, "page": 1}]


def read_csv(file_path: Path) -> List[Dict[str, Any]]:
    """CSV 解析：每行作为一个逻辑页，列以制表符拼接便于阅读。"""
    import csv

    pages: List[Dict[str, Any]] = []
    with file_path.open(encoding="utf-8", errors="ignore", newline="") as fp:
        reader = csv.reader(fp)
        for index, row in enumerate(reader, start=1):
            cells = [str(cell).strip() for cell in row]
            text = "\t".join(cell for cell in cells if cell)
            if text:
                pages.append({"text": text, "page": index})
    return pages


def read_text_file(file_path: Path) -> List[Dict[str, Any]]:
    text = file_path.read_text(encoding="utf-8", errors="ignore")
    return [{"text": text, "page": 1}]


# ---------------------------------------------------------------------------
# 统一入口
# ---------------------------------------------------------------------------
_LOADERS = {
    ".pdf": read_pdf,
    ".docx": read_docx,
    ".pptx": read_pptx,
    ".html": read_html,
    ".htm": read_html,
    ".csv": read_csv,
    ".txt": read_text_file,
    ".md": read_text_file,
}


def load_document(file_path: Path) -> List[Dict[str, Any]]:
    suffix = Path(file_path).suffix.lower()
    loader = _LOADERS.get(suffix)
    if loader is None:
        raise ValueError(f"暂不支持该文件类型：{suffix}")
    return loader(Path(file_path))


def supported_extensions() -> List[str]:
    """返回 supported 文件后缀（不带点），供 UI 文件上传限制使用。"""
    return [ext.lstrip(".") for ext in _LOADERS]
